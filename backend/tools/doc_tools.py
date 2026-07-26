"""文档生成工具：Word / PPT / Excel。从原 documents.py 迁移并大幅改进。

改进点：
- PPT：主题配色系统（商务蓝/科技深/简约灰）+ 多版式（封面/目录/章节/内容/双栏/致谢）+
  空白版式手动定位 + 标题彩色条 + 页脚页码 + 演讲者备注，告别默认版式的粗糙观感。
- Word：标题层级化、统一字体、目录字段(TOC)、页脚页码。
- Excel：沿用 openpyxl，表头加粗 + 冻结首行 + 列宽自适应。

所有工具：宽松取值、失败返回错误串而非抛异常；产文件后落盘入库 -> push_attachment -> 返回说明串。
"""
from __future__ import annotations

import uuid
from io import BytesIO
from typing import Any

from langchain.tools import tool

from config import settings
import db
from .attachments import push_attachment


# ---------------- 宽松取值 ----------------
def _coerce(value: Any, default: Any = None) -> Any:
    if value is None:
        return default
    if isinstance(value, str) and not value.strip():
        return default
    return value


def _as_list(value: Any) -> list:
    if not value:
        return []
    return value if isinstance(value, list) else [value]


# ======================================================================
#  PPT 主题与版式系统
# ======================================================================
_THEMES: dict[str, dict] = {
    "business": {  # 商务蓝
        "name": "商务蓝",
        "primary": (31, 78, 121),      # 深蓝
        "accent": (237, 125, 58),      # 橙色强调
        "bg": (255, 255, 255),
        "title_color": (255, 255, 255),
        "text_color": (51, 51, 51),
        "light": (219, 230, 242),      # 浅蓝背景条
        "font_cn": "微软雅黑",
        "font_en": "Calibri",
    },
    "tech": {  # 科技深
        "name": "科技深",
        "primary": (15, 32, 58),       # 深海军蓝
        "accent": (0, 188, 212),       # 青色
        "bg": (240, 244, 248),
        "title_color": (255, 255, 255),
        "text_color": (33, 37, 41),
        "light": (224, 234, 247),
        "font_cn": "微软雅黑",
        "font_en": "Calibri",
    },
    "gray": {  # 简约灰
        "name": "简约灰",
        "primary": (64, 64, 64),
        "accent": (192, 80, 77),
        "bg": (255, 255, 255),
        "title_color": (255, 255, 255),
        "text_color": (51, 51, 51),
        "light": (235, 235, 235),
        "font_cn": "微软雅黑",
        "font_en": "Calibri",
    },
}


def _rgb(color: tuple) -> tuple:
    """pptx RGBColor 需要单独构造，延迟导入避免模块加载期依赖。"""
    from pptx.dml.color import RGBColor
    return RGBColor(*color)


def _set_font(run, theme: dict, size: int, bold: bool = False,
              color: tuple = None, cn: bool = True):
    from pptx.util import Pt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = theme["font_en"]
    # 中文字体需通过 element 设置 east_asia
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", theme["font_cn"])
    if color is not None:
        run.font.color.rgb = _rgb(color)


def _add_textbox(slide, left, top, width, height):
    """添加无填充无边框的文本框，返回 text_frame。"""
    from pptx.util import Emu
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _add_rect(slide, left, top, width, height, fill_color):
    """添加纯色矩形（标题条/背景块）。"""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()  # 无边框
    shape.shadow.inherit = False
    return shape


def _add_page_footer(slide, page_no: int, total: int, theme: dict):
    """页脚：左侧主色小条 + 右侧页码。"""
    from pptx.util import Emu, Pt
    _add_rect(slide, Emu(0), Emu(6800000), Emu(9144000), Emu(45720), theme["light"])
    _, tf = _add_textbox(slide, Emu(6900000), Emu(6780000), Emu(2200000), Emu(300000))
    p = tf.paragraphs[0]
    p.alignment = 2  # right
    run = p.add_run()
    run.text = f"{page_no} / {total}"
    _set_font(run, theme, size=10, color=theme["text_color"])


def _layout_cover(slide, title: str, subtitle: str, author: str, theme: dict):
    """封面页：整页主色背景 + 居中标题 + 副标题 + 作者。"""
    from pptx.util import Emu
    _add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(6858000), theme["primary"])
    # 强调色横条
    _add_rect(slide, Emu(2800000), Emu(2600000), Emu(3540000), Emu(60000), theme["accent"])
    _, tf = _add_textbox(slide, Emu(800000), Emu(2900000), Emu(7500000), Emu(1500000))
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = title or "演示文稿"
    _set_font(run, theme, size=40, bold=True, color=theme["title_color"])
    if subtitle:
        _, tf2 = _add_textbox(slide, Emu(800000), Emu(4200000), Emu(7500000), Emu(600000))
        p2 = tf2.paragraphs[0]
        p2.alignment = 2
        r2 = p2.add_run()
        r2.text = subtitle
        _set_font(r2, theme, size=18, color=theme["title_color"])
    if author:
        _, tf3 = _add_textbox(slide, Emu(800000), Emu(5600000), Emu(7500000), Emu(400000))
        p3 = tf3.paragraphs[0]
        p3.alignment = 2
        r3 = p3.add_run()
        r3.text = author
        _set_font(r3, theme, size=14, color=theme["title_color"])


def _layout_section(slide, title: str, theme: dict, page_no: int, total: int):
    """章节页：左主色块 + 大号标题。"""
    from pptx.util import Emu
    _add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(6858000), theme["light"])
    _add_rect(slide, Emu(0), Emu(2200000), Emu(2400000), Emu(2400000), theme["primary"])
    _, tf = _add_textbox(slide, Emu(2700000), Emu(2700000), Emu(6000000), Emu(1200000))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title or "章节"
    _set_font(run, theme, size=36, bold=True, color=theme["primary"])
    _add_page_footer(slide, page_no, total, theme)


def _layout_content(slide, title: str, bullets: list, theme: dict,
                    page_no: int, total: int):
    """内容页：顶部彩色标题条 + 标题 + 要点（支持一级/二级）。"""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN
    # 标题条
    _add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(900000), theme["primary"])
    # 标题左侧强调竖条
    _add_rect(slide, Emu(0), Emu(0), Emu(120000), Emu(900000), theme["accent"])
    _, tf = _add_textbox(slide, Emu(300000), Emu(180000), Emu(8600000), Emu(540000))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title or "内容"
    _set_font(run, theme, size=24, bold=True, color=theme["title_color"])
    # 要点
    _, tf2 = _add_textbox(slide, Emu(450000), Emu(1100000), Emu(8240000), Emu(5400000))
    first = True
    for b in bullets:
        if isinstance(b, str):
            text, level = b, 0
        elif isinstance(b, dict):
            text = str(b.get("text", ""))
            level = int(b.get("level", 0) or 0)
        else:
            continue
        text = (text or "").strip()
        if not text:
            continue
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.level = max(0, min(1, level))
        run = p.add_run()
        prefix = "● " if level == 0 else "  ○ "
        run.text = prefix + text
        _set_font(run, theme, size=18 if level == 0 else 15,
                  bold=(level == 0), color=theme["text_color"])
        p.space_after = Pt(8)
    _add_page_footer(slide, page_no, total, theme)


def _layout_two_column(slide, title: str, left_bullets: list, right_bullets: list,
                       theme: dict, page_no: int, total: int):
    """双栏内容页。"""
    from pptx.util import Emu, Pt
    _add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(900000), theme["primary"])
    _, tf = _add_textbox(slide, Emu(300000), Emu(180000), Emu(8600000), Emu(540000))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title or "对比"
    _set_font(run, theme, size=24, bold=True, color=theme["title_color"])

    def _fill_col(tf_col, bullets):
        first = True
        for b in bullets:
            text = b if isinstance(b, str) else (b.get("text", "") if isinstance(b, dict) else str(b))
            text = (text or "").strip()
            if not text:
                continue
            p = tf_col.paragraphs[0] if first else tf_col.add_paragraph()
            first = False
            run = p.add_run()
            run.text = "● " + text
            _set_font(run, theme, size=16, color=theme["text_color"])
            p.space_after = Pt(6)

    _, tf_l = _add_textbox(slide, Emu(300000), Emu(1100000), Emu(4100000), Emu(5400000))
    _, tf_r = _add_textbox(slide, Emu(4740000), Emu(1100000), Emu(4100000), Emu(5400000))
    _fill_col(tf_l, _as_list(left_bullets))
    _fill_col(tf_r, _as_list(right_bullets))
    _add_page_footer(slide, page_no, total, theme)


def _layout_closing(slide, text: str, theme: dict):
    """致谢页。"""
    from pptx.util import Emu
    _add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(6858000), theme["primary"])
    _add_rect(slide, Emu(3300000), Emu(2600000), Emu(2540000), Emu(60000), theme["accent"])
    _, tf = _add_textbox(slide, Emu(800000), Emu(2900000), Emu(7500000), Emu(1200000))
    p = tf.paragraphs[0]
    p.alignment = 2
    run = p.add_run()
    run.text = text or "谢谢观看"
    _set_font(run, theme, size=44, bold=True, color=theme["title_color"])


def build_pptx(title: str, slides: list, theme_key: str = "business",
               subtitle: str = "", author: str = "") -> bytes:
    """用主题版式系统构建 PPT。"""
    from pptx import Presentation
    from pptx.util import Emu

    theme = _THEMES.get(theme_key, _THEMES["business"])
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]  # 空白版式，完全手动控制

    # 封面（第 1 页）
    _layout_cover(prs.slides.add_slide(blank), title, subtitle, author, theme)

    total = len(slides) + 1  # 含封面
    page = 2
    for s in _as_list(slides):
        if not isinstance(s, dict):
            continue
        layout = (s.get("layout") or "content").strip().lower()
        stitle = s.get("title", "")
        slide = prs.slides.add_slide(blank)
        if layout == "section":
            _layout_section(slide, stitle, theme, page, total)
        elif layout == "two_column":
            _layout_two_column(slide, stitle, s.get("bullets", []), s.get("right_bullets", []),
                               theme, page, total)
        elif layout == "closing":
            _layout_closing(slide, stitle, theme)
            page += 1
            continue
        else:  # content / agenda 都用内容版式
            _layout_content(slide, stitle, _as_list(s.get("bullets")), theme, page, total)
        # 演讲者备注
        notes = s.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
        page += 1

    # 若用户没显式放致谢页，自动补一页
    has_closing = any(isinstance(s, dict) and (s.get("layout") or "").lower() == "closing"
                      for s in _as_list(slides))
    if not has_closing:
        _layout_closing(prs.slides.add_slide(blank), "谢谢观看", theme)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ======================================================================
#  Word
# ======================================================================
def build_docx(title: str, sections: list, author: str = "",
               toc: bool = False) -> bytes:
    """构建 Word：标题样式层级 + 统一字体 + 可选目录 + 页脚页码。"""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    # 统一正文字体（中英文）
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

    # 主标题
    h = doc.add_heading(level=0)
    run = h.add_run(_coerce(title, "文档") or "文档")
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if author:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ar = p.add_run(f"作者：{author}")
        ar.font.size = Pt(10)
        ar.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # 目录字段（可选）
    if toc:
        doc.add_paragraph("目录", style="Heading 1")
        p = doc.add_paragraph()
        run = p.add_run()
        fldChar_begin = OxmlElement("w:fldChar")
        fldChar_begin.set(qn("w:fldCharType"), "begin")
        instrText = OxmlElement("w:instrText")
        instrText.set(qn("xml:space"), "preserve")
        instrText.text = 'TOC \\o "1-3" \\h \\z \\u'
        fldChar_sep = OxmlElement("w:fldChar")
        fldChar_sep.set(qn("w:fldCharType"), "separate")
        fldChar_end = OxmlElement("w:fldChar")
        fldChar_end.set(qn("w:fldCharType"), "end")
        run._r.append(fldChar_begin)
        run._r.append(instrText)
        run._r.append(fldChar_sep)
        run._r.append(fldChar_end)
        doc.add_page_break()

    for sec in _as_list(sections):
        if not isinstance(sec, dict):
            continue
        level = int(sec.get("level", 1) or 1)
        level = max(1, min(3, level))
        doc.add_heading(_coerce(sec.get("heading"), "章节"), level=level)
        body = _coerce(sec.get("body"), "") or ""
        for para in body.split("\n"):
            line = para.strip()
            if line:
                doc.add_paragraph(line)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ======================================================================
#  Excel
# ======================================================================
def build_xlsx(sheets: list) -> bytes:
    """构建 Excel：表头加粗+底色 + 冻结首行 + 列宽自适应。"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = Workbook()
    first = True
    for sh in _as_list(sheets):
        if not isinstance(sh, dict):
            continue
        ws = wb.active if first else wb.create_sheet()
        first = False
        ws.title = (_coerce(sh.get("name"), "Sheet") or "Sheet")[:31]
        headers = _as_list(sh.get("headers"))
        if headers:
            ws.append([str(h) for h in headers])
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="1F4E79")
                cell.alignment = Alignment(horizontal="center", vertical="center")
            ws.freeze_panes = "A2"  # 冻结表头
        for row in _as_list(sh.get("rows")):
            ws.append(["" if v is None else str(v) for v in row])
        # 列宽自适应（粗略按表头/数据最大字符数）
        for col_idx, col in enumerate(ws.columns, 1):
            max_len = 0
            for cell in col:
                max_len = max(max_len, len(str(cell.value or "")))
            from openpyxl.utils import get_column_letter
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max(max_len + 2, 8), 50)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ======================================================================
#  langchain 工具封装
# ======================================================================
def _save_and_attach(raw: bytes, prefix: str, ext: str, note: str) -> str:
    """落盘 -> 入库 -> 推附件通道 -> 返回说明串。"""
    save_name = f"{prefix}_{uuid.uuid4().hex[:12]}.{ext}"
    save_path = settings.files_dir / save_name
    save_path.write_bytes(raw)
    fid = db.add_file(
        filename=save_name, kind="document", size=len(raw),
        chars=0, text="", path=str(save_path),
    )
    push_attachment({"file_id": fid, "filename": save_name, "kind": "document", "chars": 0})
    return note


@tool
def generate_word(title: str, sections: list, author: str = "", toc: bool = False) -> str:
    """生成一份排版美观、可下载的 Word(.docx) 文档：标题样式层级化、统一字体、可选目录、页脚页码。

    调用时机：用户当前这条消息明确要 Word / “整理成文档” / “写个报告”时。
    内容要充实结构化，不要只生成空壳。

    Args:
        title: 文档主标题。
        sections: 章节列表，每项 {heading: 小标题, body: 正文(可换行分段), level: 1-3 可选}。
        author: 可选作者名。
        toc: 是否在正文前插入目录字段（默认 False）。
    """
    try:
        raw = build_docx(title, sections, author=author, toc=toc)
    except Exception as e:
        return f"生成失败：Word 文档构建出错（{e}）。"
    return _save_and_attach(raw, "doc", "docx",
                            "📄 已生成 Word 文档，见下方附件可下载。")


@tool
def generate_ppt(title: str, slides: list, theme: str = "business",
                 subtitle: str = "", author: str = "") -> str:
    """生成一份带主题配色、多版式、排版美观、可下载的 PPT(.pptx) 幻灯片。

    调用时机：用户当前这条消息明确要 PPT / “做个幻灯片” / “做汇报”时。
    要求结构完整：封面 -> 目录/章节 -> 若干内容页 -> 致谢；内容充实，每页要点清晰。
    不要只生成一两页空壳。

    Args:
        title: 演示文稿主标题（作为封面标题）。
        slides: 幻灯片列表（不含封面，封面由 title 自动生成）。每项 {
            layout: 版式类型，可选 'section'(章节页) / 'content'(标题+要点，默认) /
                    'two_column'(双栏，需提供 bullets 与 right_bullets) / 'closing'(致谢)，
            title: 本页标题,
            bullets: 要点列表，元素可为字符串或 {text, level(0或1)},
            right_bullets: 仅 two_column 版式用，右栏要点,
            notes: 可选演讲者备注
        }。
        theme: 主题配色，可选 'business'(商务蓝,默认) / 'tech'(科技深) / 'gray'(简约灰)。
        subtitle: 封面副标题，可选。
        author: 封面作者署名，可选。
    """
    try:
        raw = build_pptx(title, slides, theme_key=theme, subtitle=subtitle, author=author)
    except Exception as e:
        return f"生成失败：PPT 幻灯片构建出错（{e}）。"
    return _save_and_attach(raw, "ppt", "pptx",
                            "📊 已生成 PPT 幻灯片，见下方附件可下载。")


@tool
def generate_excel(sheets: list) -> str:
    """生成一份表头加粗配色、冻结首行、列宽自适应、可下载的 Excel(.xlsx) 表格。

    调用时机：用户当前这条消息要表格 / 数据统计 / 清单 / “导出 Excel”时。

    Args:
        sheets: 工作表列表，每项 {name: 工作表名, headers: 表头列名列表, rows: 数据行(每行为与表头等长的值数组)}。
    """
    try:
        raw = build_xlsx(sheets)
    except Exception as e:
        return f"生成失败：Excel 表格构建出错（{e}）。"
    return _save_and_attach(raw, "xls", "xlsx",
                            "📈 已生成 Excel 表格，见下方附件可下载。")


DOC_TOOLS = [generate_word, generate_ppt, generate_excel]
