'''文档自动生成：基于 Function Calling 的工具定义与 Word/PPT/Excel 打包。
模型识别到用户想要文档时返回 tool_call，后端调用本模块把结构化内容打包成
对应格式字节，落盘后作为助手消息附件下发（与图片附件共用同一条链路）。'''
from io import BytesIO
from typing import Any

# 工具引导提示词：注入到支持 FC 的模型的 system 消息里，约束调用时机与内容要求。
TOOL_GUIDE = (
    "你可以调用以下工具来生成可下载的文档文件。当用户【当前这条消息】明确要求输出 "
    "Word、PPT、Excel 或“整理成文档/做个幻灯片/导出表格”这类诉求时，调用对应的那一个工具，"
    "并把结构化、内容充实的内容作为参数传入；不要只回复文字说明。\n"
    "- generate_word：生成 Word 文档。传入标题与若干章节（每节含小标题与正文）。\n"
    "- generate_ppt：生成 PPT 幻灯片。传入标题与若干页（每页含标题与要点列表）。\n"
    "- generate_excel：生成 Excel 表格。传入若干工作表（每表含表头与数据行）。"
    "用户要表格/数据统计/清单时优先用它。\n"
    "重要约束：\n"
    "1. 一次回复最多只调用一个工具，只生成用户【当前这一条】消息所要求的那一种文档；"
    "不要因为对话历史里曾出现过别的文档请求就重复生成或补做。\n"
    "2. 平时普通问答不要调用任何工具。\n"
    "3. 调用工具后无需再重复输出文档全文，给一句简短说明即可。"
)

# OpenAI Function Calling 工具定义（JSON Schema）。仅 describe + parameters，
# 由调用方包装成 [{"type": "function", "function": {...}}] 传给 API。
_TOOLS_SCHEMA: list[dict] = [
    {
        "name": "generate_word",
        "description": "生成一份可下载的 Word(.docx) 文档。用户想要 Word 或“整理成文档”时调用。",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "文档主标题"},
                "sections": {
                    "type": "array",
                    "description": "文档章节列表",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string", "description": "章节小标题"},
                            "body": {"type": "string", "description": "章节正文，可用换行分段"},
                        },
                        "required": ["heading", "body"],
                    },
                },
            },
            "required": ["title", "sections"],
        },
    },
    {
        "name": "generate_ppt",
        "description": "生成一份可下载的 PPT(.pptx) 幻灯片。用户想要 PPT 或“做个幻灯片”时调用。",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "演示文稿主标题（作为封面）"},
                "slides": {
                    "type": "array",
                    "description": "幻灯片页列表（不含封面，封面由 title 自动生成）",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "本页标题"},
                            "bullets": {
                                "type": "array",
                                "description": "本页要点列表",
                                "items": {"type": "string"},
                            },
                        },
                        "required": ["title", "bullets"],
                    },
                },
            },
            "required": ["title", "slides"],
        },
    },
    {
        "name": "generate_excel",
        "description": "生成一份可下载的 Excel(.xlsx) 表格。用户想要 Excel 或“导出表格”时调用。",
        "parameters": {
            "type": "object",
            "properties": {
                "sheets": {
                    "type": "array",
                    "description": "工作表列表",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string", "description": "工作表名称"},
                            "headers": {
                                "type": "array",
                                "description": "表头列名",
                                "items": {"type": "string"},
                            },
                            "rows": {
                                "type": "array",
                                "description": "数据行，每行为与表头等长的值数组",
                                "items": {"type": "array", "items": {}},
                            },
                        },
                        "required": ["name", "headers", "rows"],
                    },
                },
            },
            "required": ["sheets"],
        },
    },
]


def tools_for_api() -> list[dict]:
    """包装成 OpenAI API 要求的 tools 格式。"""
    return [
        {"type": "function", "function": t} for t in _TOOLS_SCHEMA
    ]


# ---------------- 文档打包 ----------------
def _coerce(value: Any, default: Any = None) -> Any:
    """宽松取值：None 或非预期类型时给默认。"""
    if value is None:
        return default
    if isinstance(value, str) and not value.strip():
        return default
    return value


def build_docx(title: str, sections: list[dict]) -> bytes:
    """把 标题 + 章节 列表打包成 .docx 字节。"""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(_coerce(title, "文档") or "文档", level=0)
    for sec in sections or []:
        if not isinstance(sec, dict):
            continue
        doc.add_heading(_coerce(sec.get("heading"), "章节"), level=1)
        body = _coerce(sec.get("body"), "") or ""
        for para in body.split("\n"):
            line = para.strip()
            if line:
                doc.add_paragraph(line)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pptx(title: str, slides: list[dict]) -> bytes:
    """把 标题 + 幻灯片列表 打包成 .pptx 字节。封面页用 title，其余逐页渲染。"""
    from pptx import Presentation

    prs = Presentation()
    # 封面页
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = _coerce(title, "演示文稿") or "演示文稿"
    for slide in slides or []:
        if not isinstance(slide, dict):
            continue
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = _coerce(slide.get("title"), "标题") or "标题"
        body = s.placeholders[1]
        tf = body.text_frame
        bullets = slide.get("bullets") or []
        if not bullets:
            tf.text = ""
            continue
        tf.text = str(bullets[0])
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = str(b)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_xlsx(sheets: list[dict]) -> bytes:
    """把若干工作表打包成 .xlsx 字节。表头加粗，数据行逐行写入。"""
    from openpyxl import Workbook

    wb = Workbook()
    # 默认首表会被创建，先复用或删掉
    first = True
    for sh in sheets or []:
        if not isinstance(sh, dict):
            continue
        ws = wb.active if first else wb.create_sheet()
        first = False
        ws.title = (_coerce(sh.get("name"), "Sheet") or "Sheet")[:31]
        headers = sh.get("headers") or []
        if headers:
            ws.append([str(h) for h in headers])
            for cell in ws[1]:
                cell.font = cell.font.copy(bold=True)
        for row in sh.get("rows") or []:
            ws.append(["" if v is None else str(v) for v in row])
    if first:
        # 一个工作表都没有：保留默认空表
        pass
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------- 工具分发 ----------------
# 工具名 -> (打包函数, 输出扩展名, 前端说明文案)。返回字节供后端落盘。
_DISPATCH = {
    "generate_word": (build_docx, "docx", "📄 已生成 Word 文档，见下方附件可下载。"),
    "generate_ppt": (build_pptx, "pptx", "📊 已生成 PPT 幻灯片，见下方附件可下载。"),
    "generate_excel": (build_xlsx, "xlsx", "📈 已生成 Excel 表格，见下方附件可下载。"),
}


def dispatch_tool(name: str, args: dict) -> tuple[bytes, str, str]:
    """按工具名路由到对应打包函数，返回 (文件字节, 扩展名, 前端说明文案)。
    未知工具名抛 ValueError，由调用方转成 error 事件。"""
    if name not in _DISPATCH:
        raise ValueError(f"未知的文档生成工具：{name}")
    builder, ext, note = _DISPATCH[name]
    if name == "generate_word":
        raw = builder(args.get("title"), args.get("sections"))
    elif name == "generate_ppt":
        raw = builder(args.get("title"), args.get("slides"))
    else:  # generate_excel
        raw = builder(args.get("sheets"))
    return raw, ext, note
